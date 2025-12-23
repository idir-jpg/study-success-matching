import os
from pptx import Presentation
import pandas as pd
from email.message import EmailMessage
from email.utils import make_msgid

def update_ppt(selected_row, excel_file_path, pptx_file_path, pdf_file_folder, oft_template_path=None):
    df = pd.read_excel(excel_file_path, sheet_name=None)
    selected_id = selected_row['Id']

    details_df = df['Profils_élèves'].loc[df['Profils_élèves']['id'] == selected_id]
    if details_df.empty:
        raise ValueError(f"Aucun élève trouvé avec l'ID {selected_id}.")
    details = details_df.iloc[0]

    suivi_info = df['Suivi'].loc[df['Suivi']['Id'] == selected_id].iloc[0]
    to_email = suivi_info['Mail']

    prs = Presentation(pptx_file_path)

    def set_shape_text(shape, new_text):
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = ''
                if paragraph.runs:
                    paragraph.runs[0].text = new_text
                else:
                    paragraph.text = new_text

    set_shape_text(prs.slides[0].shapes[10], f"{details['Prénom']} {details['Nom'].upper()}")

    visuel_verbal = "Visuel" if details['Visuel'] > details['Verbal'] else "Verbal"
    sensoriel_intuitif = "Sensoriel" if details['Sensoriel'] > details['Intuitif'] else "Intuitif"
    actif_reflexif = "Actif" if details['Actif'] > details['Réflexif'] else "Réflexif"
    sequentiel_global = "Séquentiel" if details['Séquentiel'] > details['Global'] else "Global"

    set_shape_text(prs.slides[0].shapes[1], visuel_verbal)
    set_shape_text(prs.slides[0].shapes[2], sensoriel_intuitif)
    set_shape_text(prs.slides[0].shapes[3], actif_reflexif)
    set_shape_text(prs.slides[0].shapes[7], sequentiel_global)

    descriptions = {
        "Visuel": "Ton super-pouvoir, c’est la mémoire des images ! Schémas, mindmaps, vidéos, couleurs…",
        "Verbal": "Si c’est expliqué à l’oral ou à l’écrit, tu captes vite !",
        "Sensoriel": "Tu as un esprit logique et concret. Tu n’aimes pas les imprévus...",
        "Intuitif": "La routine t’ennuie ! Tu aimes découvrir de nouvelles idées...",
        "Actif": "Tu apprends en faisant : expérimenter, manipuler, discuter...",
        "Réflexif": "Tu aimes prendre ton temps pour comprendre en profondeur.",
        "Séquentiel": "Tu préfères apprendre étape par étape, en suivant une logique.",
        "Global": "Tu as besoin de comprendre la vision d’ensemble avant les détails."
    }

    set_shape_text(prs.slides[0].shapes[4], descriptions[visuel_verbal])
    set_shape_text(prs.slides[0].shapes[5], descriptions[sensoriel_intuitif])
    set_shape_text(prs.slides[0].shapes[6], descriptions[sequentiel_global])
    set_shape_text(prs.slides[0].shapes[9], descriptions[actif_reflexif])

    # Enregistrement du PowerPoint mis à jour
    new_pptx_path = os.path.join(pdf_file_folder, f"{details['Prénom']}_{details['Nom']}_profil.pptx")
    prs.save(new_pptx_path)

    # === Génération de l'e-mail HTML avec image intégrée ===
    image_path = os.path.join(os.path.dirname(__file__), "etapes.png")
    image_cid = make_msgid(domain="studysuccess.local")[1:-1]

    mail_html = f"""
    <html>
      <body>
        Bonjour,<br><br>
        Merci d'avoir complété le test de profil d'apprentissage de Study Success, voici les résultats :<br><br>
        Voici les prochaines étapes :<br><br>
        <img src="cid:{image_cid}"><br><br>
        Merci pour votre confiance,<br>
        Excellente journée,
      </body>
    </html>
    """

    msg = EmailMessage()
    msg['Subject'] = f"Résultats test de profil de {details['Prénom']}"
    msg['To'] = to_email
    msg.set_content("Ce message contient un contenu HTML.")
    msg.add_alternative(mail_html, subtype='html')

    # Pièce jointe : PPTX
    with open(new_pptx_path, 'rb') as f:
        msg.add_attachment(
            f.read(),
            maintype='application',
            subtype='vnd.openxmlformats-officedocument.presentationml.presentation',
            filename=os.path.basename(new_pptx_path)
        )

    # Image intégrée
    if os.path.exists(image_path):
        with open(image_path, 'rb') as img:
            msg.get_payload()[1].add_related(
                img.read(),
                maintype='image',
                subtype='png',
                cid=f"<{image_cid}>"
            )

    # Sauvegarde .emltpl
    eml_path = os.path.join(pdf_file_folder, f"profil_{details['Prénom']}_{details['Nom']}.emltpl")
    with open(eml_path, 'wb') as f:
        f.write(bytes(msg))

    return eml_path
